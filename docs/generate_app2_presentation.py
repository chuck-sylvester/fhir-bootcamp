#!/usr/bin/env python3
"""
generate_app2_presentation.py
Generates docs/app2-tutorial.pptx — a full technical tutorial presentation
for the fhir-bootcamp app2 application.

Run from the project root with the venv active:
    python docs/generate_app2_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette ──────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
BLUE    = RGBColor(0x00, 0x78, 0xB4)
TEAL    = RGBColor(0x00, 0x96, 0x99)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
L_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
M_GRAY  = RGBColor(0x77, 0x88, 0x99)
D_GRAY  = RGBColor(0x22, 0x33, 0x44)
MUTED   = RGBColor(0xBB, 0xCC, 0xDD)
ORANGE  = RGBColor(0xE8, 0x7A, 0x1E)
GREEN   = RGBColor(0x2E, 0x86, 0x48)
CODE_BG = RGBColor(0x1A, 0x1A, 0x2E)
CODE_FG = RGBColor(0xCB, 0xD7, 0xE4)

# ── Slide geometry ─────────────────────────────────────────────────────────────
SW = Inches(13.33)   # slide width  (16:9 widescreen)
SH = Inches(7.5)     # slide height
BAR_H  = Inches(1.25)
MARGIN = Inches(0.55)
CT     = BAR_H + Inches(0.3)          # content top
CW     = SW - 2 * MARGIN              # content width
CH     = SH - CT - Inches(0.25)       # content height

# ── Low-level helpers ──────────────────────────────────────────────────────────

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return sl

def box(sl, x, y, w, h, fill, no_line=True):
    sh = sl.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if no_line:
        sh.line.color.rgb = fill
    return sh

def tb(sl, text, x, y, w, h, size=18, color=D_GRAY,
       bold=False, italic=False, align=PP_ALIGN.LEFT,
       font="Calibri", wrap=True):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return txb

def blist(sl, items, x, y, w, h, base=19, spacing=5):
    """
    items: list of (text, level) tuples.  level 0 = bullet, 1 = sub-bullet.
    """
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sz   = base if lvl == 0 else base - 3
        col  = D_GRAY if lvl == 0 else M_GRAY
        bul  = "▸  " if lvl == 0 else "    –  "
        p.space_before = Pt(spacing if lvl == 0 else 2)
        r = p.add_run()
        r.text = bul + text
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.name = "Calibri"
    return txb

def code(sl, src, x, y, w, h, cap=None):
    box(sl, x, y, w, h, CODE_BG, no_line=False)
    pad = Inches(0.22)
    txb = sl.shapes.add_textbox(x + pad, y + pad, w - 2*pad, h - 2*pad)
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(src.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.color.rgb = CODE_FG
        r.font.name = "Courier New"
    if cap:
        tb(sl, cap, x, y + h + Inches(0.08), w, Inches(0.3),
           size=12, color=M_GRAY, italic=True)

# ── Standard slide chrome ──────────────────────────────────────────────────────

def title_bar(sl, title, sub=None):
    box(sl, 0, 0, SW, BAR_H, NAVY)
    box(sl, 0, 0, Inches(0.14), BAR_H, TEAL)
    if sub:
        tb(sl, title, MARGIN, Inches(0.15), CW, Inches(0.7),
           size=24, color=WHITE, bold=True)
        tb(sl, sub, MARGIN, Inches(0.82), CW, Inches(0.38),
           size=13, color=MUTED, italic=True)
    else:
        tb(sl, title, MARGIN, Inches(0.22), CW, Inches(0.9),
           size=27, color=WHITE, bold=True)

# ── Slide type constructors ────────────────────────────────────────────────────

def slide_title(prs, title, line2, line3=""):
    sl = new_slide(prs)
    box(sl, 0, 0, SW, SH, NAVY)
    box(sl, 0, 0, Inches(0.18), SH, TEAL)
    box(sl, 0, SH - Inches(0.08), SW, Inches(0.08), ORANGE)
    tb(sl, title, Inches(0.6), Inches(1.6), Inches(12), Inches(1.6),
       size=38, color=WHITE, bold=True, wrap=True)
    tb(sl, line2, Inches(0.6), Inches(3.4), Inches(12), Inches(0.7),
       size=22, color=TEAL, bold=True)
    if line3:
        tb(sl, line3, Inches(0.6), Inches(4.15), Inches(12), Inches(0.5),
           size=16, color=MUTED)

def slide_section(prs, num, title, desc=""):
    sl = new_slide(prs)
    box(sl, 0, 0, SW, SH, BLUE)
    box(sl, 0, 0, Inches(0.18), SH, TEAL)
    box(sl, 0, SH - Inches(0.08), SW, Inches(0.08), ORANGE)
    tb(sl, f"SECTION {num}", Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
       size=16, color=MUTED, bold=True, font="Calibri")
    tb(sl, title, Inches(0.6), Inches(2.0), Inches(12), Inches(1.8),
       size=36, color=WHITE, bold=True, wrap=True)
    if desc:
        tb(sl, desc, Inches(0.6), Inches(4.2), Inches(11), Inches(0.7),
           size=16, color=MUTED, italic=True)

def slide_bullets(prs, title, items, sub=None):
    sl = new_slide(prs)
    title_bar(sl, title, sub)
    blist(sl, items, MARGIN, CT, CW, CH)
    return sl

def slide_code(prs, title, src, cap=None, sub=None, extra_bullets=None):
    sl = new_slide(prs)
    title_bar(sl, title, sub)
    if extra_bullets:
        blist(sl, extra_bullets, MARGIN, CT, CW * 0.42, CH)
        code(sl, src, MARGIN + CW * 0.45, CT, CW * 0.55, CH - Inches(0.1), cap)
    else:
        code(sl, src, MARGIN, CT, CW, CH - Inches(0.1), cap)
    return sl

def slide_two_col(prs, title, lh, li, rh, ri, sub=None):
    """Two-column comparison slide."""
    sl = new_slide(prs)
    title_bar(sl, title, sub)
    col_w = CW * 0.47
    gap   = CW * 0.06
    # Left column header
    box(sl, MARGIN, CT, col_w, Inches(0.45), NAVY)
    tb(sl, lh, MARGIN + Inches(0.12), CT + Inches(0.05), col_w - Inches(0.2),
       Inches(0.38), size=14, color=WHITE, bold=True)
    blist(sl, li, MARGIN, CT + Inches(0.55), col_w, CH - Inches(0.55), base=17)
    # Right column header
    rx = MARGIN + col_w + gap
    box(sl, rx, CT, col_w, Inches(0.45), TEAL)
    tb(sl, rh, rx + Inches(0.12), CT + Inches(0.05), col_w - Inches(0.2),
       Inches(0.38), size=14, color=WHITE, bold=True)
    blist(sl, ri, rx, CT + Inches(0.55), col_w, CH - Inches(0.55), base=17)
    return sl

def slide_closing(prs, title, line2=""):
    sl = new_slide(prs)
    box(sl, 0, 0, SW, SH, NAVY)
    box(sl, 0, 0, Inches(0.18), SH, TEAL)
    box(sl, 0, SH - Inches(0.08), SW, Inches(0.08), ORANGE)
    tb(sl, title, Inches(0.6), Inches(2.2), Inches(12), Inches(1.5),
       size=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if line2:
        tb(sl, line2, Inches(0.6), Inches(4.0), Inches(12), Inches(0.7),
           size=20, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# ══════════════════════════════════════════════════════════════════════════════

def gen_title(prs):
    slide_title(prs,
        "Building a SMART on FHIR Application",
        "FastAPI  ·  HTMX  ·  Jinja2  ·  HTTPX",
        "A Full Technical Tutorial  |  fhir-bootcamp / app2")

def gen_s1_intro(prs):
    slide_section(prs, 1, "Introduction",
        "What we're building, who this is for, and how to follow along")

    slide_bullets(prs, "What We're Building", [
        ("A server-side rendered (SSR) web application in Python", 0),
        ("Authenticates users via Epic Sandbox — SMART on FHIR OAuth 2.0", 0),
        ("Calls the Epic FHIR R4 API for real patient data", 0),
        ("No JavaScript framework — HTMX handles interactivity", 0),
        ("", 0),
        ("Four live FHIR queries from the patient portal:", 0),
        ("GET /Patient — demographics", 1),
        ("GET /MedicationRequest — active prescriptions", 1),
        ("GET /Observation?category=laboratory — lab results", 1),
        ("GET /Observation?category=vital-signs — vitals", 1),
        ("", 0),
        ("Plus an OIDC ID Token dialog showing decoded identity claims", 0),
    ])

    slide_bullets(prs, "Prerequisites", [
        ("Python — comfortable with async/await basics", 0),
        ("Web development — some FastAPI or Flask experience helpful", 0),
        ("HTTP fundamentals — GET/POST, headers, status codes, cookies", 0),
        ("HTML/CSS — able to read and write basic markup", 0),
        ("", 0),
        ("No FHIR or OAuth experience required", 0),
        ("All concepts are introduced from the ground up", 1),
        ("No healthcare domain knowledge assumed", 1),
    ])

    slide_bullets(prs, "Repository Structure", [
        ("fhir-bootcamp/ — project root", 0),
        (".venv/  requirements.txt  .env — shared across all apps", 1),
        ("app2/ — this tutorial's application", 0),
        ("main.py  config.py  __init__.py", 1),
        ("routers/auth.py   — OAuth routes", 1),
        ("routers/pages.py  — UI and FHIR routes", 1),
        ("templates/  static/  services/", 1),
        ("docs/app2-guide.md — companion developer guide", 0),
    ])

    slide_bullets(prs, "How to Follow Along", [
        ("1.  Clone the repository", 0),
        ("2.  Create a virtual environment and install dependencies", 0),
        ("     python3.12 -m venv .venv && source .venv/bin/activate", 1),
        ("     pip install -r requirements.txt", 1),
        ("3.  Register a free developer app at open.epic.com", 0),
        ("4.  Copy your Client ID and Client Secret into .env", 0),
        ("5.  Run the application from the project root:", 0),
        ("     uvicorn app2.main:app --reload --port 8000", 1),
        ("6.  Open http://localhost:8000 in a browser", 0),
        ("", 0),
        ("app2-guide.md documents every step in full detail", 0),
    ])


def gen_s2_fhir(prs):
    slide_section(prs, 2, "FHIR Fundamentals",
        "Resources, REST conventions, and SMART on FHIR")

    slide_bullets(prs, "What is FHIR?", [
        ("Fast Healthcare Interoperability Resources — HL7 FHIR R4", 0),
        ("A REST API specification for exchanging healthcare data", 0),
        ("R4 (Release 4) is the current stable version and the standard for US EHRs", 0),
        ("", 0),
        ("FHIR became legally required for US EHR vendors in 2021", 0),
        ("21st Century Cures Act mandates FHIR R4 patient data access APIs", 1),
        ("", 0),
        ("Implemented by all major US EHR vendors:", 0),
        ("Epic, Oracle Cerner, Microsoft Azure Health APIs, Google Cloud Healthcare API", 1),
        ("", 0),
        ("Once you know FHIR, you can query any compliant system with the same API", 0),
    ])

    slide_bullets(prs, "FHIR Resources: The Data Model", [
        ("In FHIR, every piece of clinical information is a 'Resource'", 0),
        ("Resources are structured JSON objects with a consistent shape", 0),
        ("", 0),
        ("Resources used in app2:", 0),
        ("Patient        — name, date of birth, gender, address", 1),
        ("MedicationRequest  — prescriptions and medication orders", 1),
        ("Observation    — lab results, vital signs, questionnaire responses", 1),
        ("", 0),
        ("Other common resources (not in app2 yet):", 0),
        ("Condition — diagnoses  |  Encounter — visits  |  Practitioner — clinicians", 1),
        ("AllergyIntolerance  |  Procedure  |  Immunization  |  DiagnosticReport", 1),
    ])

    slide_bullets(prs, "FHIR R4 REST API Conventions", [
        ("Base URL (Epic sandbox):", 0),
        ("https://fhir.epic.com/interconnect-fhir-oauth/api/FHIR/R4", 1),
        ("", 0),
        ("Read a specific resource by ID:", 0),
        ("GET /Patient/{id}  →  returns one Patient resource", 1),
        ("", 0),
        ("Search with parameters:", 0),
        ("GET /MedicationRequest?patient={id}  →  returns a Bundle", 1),
        ("GET /Observation?patient={id}&category=laboratory  →  Bundle", 1),
        ("", 0),
        ("Response format:", 0),
        ("Single resource read  →  the resource JSON directly", 1),
        ("Search  →  a Bundle containing entry[] of matching resources", 1),
        ("Content-Type: application/fhir+json  (preferred over application/json)", 1),
    ])

    slide_bullets(prs, "SMART on FHIR", [
        ("SMART = Substitutable Medical Apps Reusable Technologies", 0),
        ("Adds OAuth 2.0 + OpenID Connect (OIDC) on top of FHIR", 0),
        ("The standard authorization layer for healthcare app integration", 0),
        ("", 0),
        ("What SMART provides:", 0),
        ("A defined OAuth 2.0 flow for clinical apps (Authorization Code)", 1),
        ("Scopes that map to specific FHIR resource types and permissions", 1),
        ("Identity (OIDC) so the app knows which patient or clinician logged in", 1),
        ("Launch context — the in-scope patient ID travels with the token", 1),
        ("", 0),
        ("Scope format:  patient/Patient.read  patient/Observation.read", 0),
        ("Epic uses abbreviated notation in token response: .r for read, .s for search", 1),
    ])

    slide_bullets(prs, "Epic Sandbox Overview", [
        ("Non-production sandbox environment — free for developers", 0),
        ("Register at open.epic.com — no institutional affiliation required", 0),
        ("", 0),
        ("What the sandbox provides:", 0),
        ("Pre-populated test patients (e.g., Camila Lopez)", 1),
        ("A fully functional OAuth 2.0 + SMART on FHIR flow", 1),
        ("Real FHIR R4 API responses using synthetic patient data", 1),
        ("", 0),
        ("Key differences from production:", 0),
        ("No approval process — sandbox apps provision in minutes", 1),
        ("Scope behavior differs (covered in the App Registration section)", 1),
        ("Test patients have limited but representative clinical histories", 1),
    ])

    slide_bullets(prs, "FHIR Resources Used in app2", [
        ("Patient — GET /Patient", 0),
        ("Uses a 'me' context: automatically scoped to the logged-in patient", 1),
        ("No search parameters needed — Epic infers the patient from the token", 1),
        ("", 0),
        ("MedicationRequest — GET /MedicationRequest?patient=<id>", 0),
        ("Requires explicit patient ID search parameter", 1),
        ("Returns a Bundle of prescriptions and medication orders", 1),
        ("", 0),
        ("Observation — GET /Observation?patient=<id>&category=<code>", 0),
        ("category=laboratory  →  CBC, metabolic panels, urinalysis, cultures", 1),
        ("category=vital-signs  →  BP, heart rate, temperature, SpO2, height, weight", 1),
        ("Both require explicit patient ID and category filter", 1),
    ])


def gen_s3_arch(prs):
    slide_section(prs, 3, "Application Architecture",
        "SSR vs SPA, the full request lifecycle, and project structure")

    slide_two_col(prs,
        "Server-Side Rendering vs Single-Page Application",
        "SSR + HTMX (app2)",
        [
            ("HTML rendered on the server in Python", 0),
            ("Browser receives finished HTML", 0),
            ("HTMX injects fragments — no reload", 0),
            ("Logic lives in Python (typed, testable)", 0),
            ("No JavaScript framework or build pipeline", 0),
            ("One runtime to reason about", 0),
            ("Ideal when data and auth live server-side", 0),
        ],
        "Single-Page Application (React / Vue)",
        [
            ("HTML built by JavaScript in the browser", 0),
            ("Browser receives JSON data + JS bundle", 0),
            ("Client-side routing — no reload", 0),
            ("Logic split across Python + JavaScript", 0),
            ("npm, webpack, TypeScript compiler required", 0),
            ("Two runtimes: Python backend + JS frontend", 0),
            ("Ideal for rich local interactivity", 0),
        ]
    )

    slide_bullets(prs, "The Full Tech Stack", [
        ("FastAPI + uvicorn  — async Python web framework (ASGI)", 0),
        ("HTMX  — server-driven UI interactions without a JS framework", 0),
        ("Jinja2  — server-side HTML templating with template inheritance", 0),
        ("HTTPX  — async HTTP client for outbound calls to Epic FHIR API", 0),
        ("Starlette SessionMiddleware  — signed session cookies", 0),
        ("itsdangerous  — cookie signing (via TimestampSigner)", 0),
        ("Pydantic Settings  — typed configuration from .env", 0),
        ("", 0),
        ("Everything runs in a single Python process.", 0),
        ("No message queues, no microservices, no JavaScript build step.", 0),
    ])

    slide_bullets(prs, "Request Lifecycle: A Button Click to DOM Update", [
        ("1.  User clicks 'GET /Patient' in the browser", 0),
        ("     HTMX intercepts the click — sends GET /fhir/patient with session cookie", 1),
        ("2.  SessionMiddleware decodes and verifies the signed cookie", 0),
        ("     Populates request.session dict before the handler runs", 1),
        ("3.  Router matches /fhir/patient → fhir_get_patient() handler", 0),
        ("4.  Handler reads session_id, looks up access_token in app.state.token_store", 0),
        ("5.  Handler awaits HTTP GET to Epic FHIR with Authorization: Bearer <token>", 0),
        ("     httpx.AsyncClient — non-blocking, event loop serves other requests", 1),
        ("6.  Epic validates JWT, checks scope, returns FHIR Bundle JSON", 0),
        ("7.  Handler formats JSON as <pre> HTML fragment, returns HTMLResponse", 0),
        ("8.  HTMX injects the fragment into #result — no navigation, no reload", 0),
    ])

    slide_bullets(prs, "Why SSR + HTMX for a FHIR App?", [
        ("Access tokens never reach the browser", 0),
        ("Bearer tokens are kept server-side; the browser holds only a session ID", 1),
        ("HTMX fragment requests are same-origin — session cookie sent automatically", 1),
        ("", 0),
        ("Python is where FHIR logic belongs", 0),
        ("Parsing Bundles, checking scopes, decoding JWTs — all in typed Python", 1),
        ("No JSON-to-JavaScript-to-HTML transformation pipeline", 1),
        ("", 0),
        ("Simpler deployment and operations", 0),
        ("One process: no separate API server + frontend build", 1),
        ("Easier to reason about security — single trust boundary", 1),
    ])


def gen_s4_stack(prs):
    slide_section(prs, 4, "Tech Stack Deep Dive",
        "FastAPI · HTMX · Jinja2 · HTTPX")

    # FastAPI
    slide_bullets(prs, "FastAPI: Async-First Web Framework", [
        ("Built on Starlette (ASGI) and Pydantic", 0),
        ("ASGI = Asynchronous Server Gateway Interface — the modern async Python web standard", 1),
        ("uvicorn is the ASGI server that runs the FastAPI app", 1),
        ("", 0),
        ("Type-hint based — route parameters and bodies declared with Python types", 0),
        ("Auto-generates OpenAPI docs at /docs  (Swagger UI)", 0),
        ("", 0),
        ("Async-first: designed for I/O-bound workloads like HTTP API calls", 0),
        ("A FastAPI process can handle many concurrent requests efficiently", 1),
        ("The event loop suspends while waiting for network I/O (Epic FHIR calls)", 1),
    ])

    slide_code(prs,
        "FastAPI: Async / Await",
        """\
# Every route handler is async — it can yield control during I/O
async def fhir_get_patient(request: Request):

    # The await keyword suspends this handler during the network call.
    # While Epic is responding (~100-500ms), the event loop serves
    # other incoming requests.
    response = await request.app.state.http_client.get(
        f"{settings.epic_fhir_base_url}/Patient",
        headers={"Authorization": f"Bearer {access_token}"},
    )
    return HTMLResponse(content=f'<pre>{response.text}</pre>')

# Rule: async propagates upward.
# If a function contains await, it must be async def.
# If a function calls an async def, it must itself be async def.""",
        "Never use synchronous blocking calls (requests.get, time.sleep) inside an async handler"
    )

    slide_bullets(prs, "FastAPI: APIRouter and Application Assembly", [
        ("main.py is intentionally thin — it creates the app and mounts routers", 0),
        ("Route handlers live in routers/auth.py and routers/pages.py", 0),
        ("", 0),
        ("Each router is an independent module with no knowledge of the others:", 0),
        ("auth.py   — /auth/login, /auth/callback, /auth/logout", 1),
        ("pages.py  — /, /patient, /fhir/patient, /fhir/medication, /fhir/labreport, /fhir/vitalsigns", 1),
        ("", 0),
        ("Routers are mounted in main.py:", 0),
        ("app.include_router(auth_router, prefix='/auth')", 1),
        ("/login in auth.py  →  /auth/login in the full app", 1),
        ("", 0),
        ("include_in_schema=False — excludes a route from /docs (used for HTMX endpoints)", 0),
    ])

    slide_bullets(prs, "FastAPI: The Request Object", [
        ("Every route handler receives a Request — the integration point for everything:", 0),
        ("", 0),
        ("request.session          — decoded session cookie payload (dict-like)", 0),
        ("request.app              — the FastAPI application instance", 0),
        ("request.app.state        — app.state namespace (token_store, http_client...)", 0),
        ("request.query_params     — URL query string (?code=... from Epic callback)", 0),
        ("request.headers          — incoming HTTP headers", 0),
        ("request.url_for('name')  — reverse URL lookup, generates absolute URLs", 0),
        ("", 0),
        ("Why request.app instead of importing app directly?", 0),
        ("Avoids circular imports between routers and main.py", 1),
        ("Any router can access app state without importing from any other module", 1),
    ])

    slide_bullets(prs, "FastAPI: Middleware Pipeline", [
        ("Middleware wraps every request/response cycle automatically", 0),
        ("Registered once in main.py — applies to all routes", 0),
        ("", 0),
        ("SessionMiddleware (from Starlette):", 0),
        ("Inbound: reads session cookie → verifies itsdangerous HMAC signature", 1),
        ("         decodes JSON payload → attaches as request.session dict", 1),
        ("Outbound: serializes request.session → signs → sets Set-Cookie header", 1),
        ("Any mutation to request.session is automatically persisted — no explicit save", 1),
        ("", 0),
        ("Key configuration:", 0),
        ("https_only=False  — must be True in production (requires HTTPS)", 1),
        ("same_site='lax'   — allows cookie on top-level cross-site GET (Epic callback)", 1),
    ])

    # HTMX
    slide_bullets(prs, "HTMX: The Hypermedia Model", [
        ("Core idea: HTML with links and forms already describes what a user can do.", 0),
        ("HTMX extends this by letting any HTML element make HTTP requests.", 0),
        ("", 0),
        ("The critical rule: HTMX endpoints return HTML, not JSON.", 0),
        ("HTMX takes the server's response and injects it directly into the DOM.", 1),
        ("There is no JavaScript step that reads JSON and builds markup.", 1),
        ("", 0),
        ("In app2, every FHIR route returns an HTML fragment:", 0),
        ("Success:   <pre class='fhir-json'>{formatted_json}</pre>", 1),
        ("Error:     <p class='fhir-error'>Message</p>", 1),
        ("", 0),
        ("HTMX is loaded from a single CDN script tag — no build step, no npm.", 0),
    ])

    slide_bullets(prs, "HTMX: Attribute Vocabulary", [
        ("hx-get / hx-post / hx-put / hx-delete  — HTTP method and URL", 0),
        ("hx-target    — CSS selector of the element to update", 0),
        ("hx-swap      — how to place the response (innerHTML, outerHTML, beforeend...)", 0),
        ("hx-indicator — CSS selector of a loading element to show during request", 0),
        ("hx-disabled-elt  — disable this element while request is in flight", 0),
        ("", 0),
        ("Additional attributes worth knowing:", 0),
        ("hx-trigger   — what fires the request (click, change, keyup delay:500ms)", 1),
        ("hx-push-url  — update the browser address bar without navigation", 1),
        ("hx-boost     — upgrade all <a> and <form> tags in a subtree to HTMX", 1),
        ("hx-confirm   — show a confirm dialog before issuing the request", 1),
        ("", 0),
        ("hx-disabled-elt='this'  is essential on action buttons — prevents double-clicks", 0),
    ])

    slide_bullets(prs, "HTMX: Error Handling", [
        ("HTMX does NOT swap the response body when the server returns 4xx or 5xx.", 0),
        ("It fires htmx:responseError and leaves the page unchanged.", 0),
        ("", 0),
        ("The correct pattern for HTMX error responses:", 0),
        ("Return HTTP 200 with a styled error HTML fragment", 1),
        ("HTMX injects the error message into #result — in context, no page reload", 1),
        ("", 0),
        ("app2 uses a _error_html() helper for all error paths:", 0),
        ("<p class='fhir-error'>Message</p>", 1),
        ("<pre class='fhir-error-detail'>raw Epic response body</pre>  (optional)", 1),
        ("", 0),
        ("Every FHIR route handler returns error HTML at 200 — never HTTPException", 0),
        ("HTTPException is reserved for auth routes where a full error page is appropriate", 0),
    ])

    # Jinja2
    slide_bullets(prs, "Jinja2: Server-Side Templating", [
        ("Templates are HTML files with Jinja2 placeholders — rendered on the server", 0),
        ("{{ variable }}   — render a context variable", 0),
        ("{% if %}  {% for %}  {% block %}  — control flow and structure", 0),
        ("", 0),
        ("Context dict: the only bridge from Python to the template", 0),
        ("Every variable the template uses must be in this dict", 1),
        ("Pass pre-formatted values — don't make templates compute or format data", 1),
        ("", 0),
        ("Autoescaping: Jinja2 escapes < > & \" ' by default", 0),
        ("This is the primary built-in XSS defense — user input becomes safe text", 1),
        ("The | safe filter bypasses escaping — use only for trusted server-generated content", 1),
        ("Never apply | safe to user-supplied input", 1),
    ])

    slide_bullets(prs, "Jinja2: Template Inheritance", [
        ("Base template defines layout and declares named override points (blocks)", 0),
        ("Child templates fill those blocks with their specific content", 0),
        ("", 0),
        ("app2's portal hierarchy:", 0),
        ("portal.html  —  base layout: HTMX script, CSS, Font Awesome, page structure", 1),
        ("  {% block actions %}  — FHIR action buttons", 1),
        ("  {% block content %}  — #result div for HTMX injection", 1),
        ("patient.html  —  extends portal.html, fills both blocks", 1),
        ("", 0),
        ("Use {# Jinja2 comment #} in base templates, not <!-- HTML comment -->", 0),
        ("Jinja2 evaluates {% block %} tags even inside HTML comments → TemplateSyntaxError", 1),
        ("Jinja2 block comments {# ... #} strip content before any processing", 1),
        ("", 0),
        ("{{ super() }}  renders parent block content — use to append, not replace", 0),
    ])

    # HTTPX
    slide_bullets(prs, "HTTPX: Async HTTP Client", [
        ("httpx.AsyncClient is the async equivalent of the requests library", 0),
        ("Used for all outbound HTTP calls: token exchange + all FHIR API requests", 0),
        ("", 0),
        ("Why not requests?", 0),
        ("requests is synchronous — blocks the calling thread during network I/O", 1),
        ("In an async handler, blocking the thread = blocking the entire event loop", 1),
        ("httpx.AsyncClient suspends with await — event loop serves other requests", 1),
        ("", 0),
        ("Shared AsyncClient via lifespan (app.state.http_client):", 0),
        ("Created once at startup — maintains a pool of persistent TCP connections", 1),
        ("Epic TLS handshake overhead paid once, not on every request", 1),
        (".aclose() on shutdown — drains in-flight requests, releases socket FDs", 1),
        ("", 0),
        ("Always check response.is_success before calling response.json()", 0),
        ("Error responses from Epic may return non-JSON bodies → JSONDecodeError", 1),
    ])


def gen_s5_reg(prs):
    slide_section(prs, 5, "Epic App Registration",
        "Confidential clients, API capabilities, and scope behavior")

    slide_two_col(prs,
        "Confidential vs Public OAuth Clients",
        "Confidential Client (app2)",
        [
            ("Server-side application", 0),
            ("Client secret stored securely on server", 0),
            ("Secret never reaches the browser", 0),
            ("Uses client_secret_post in token exchange", 0),
            ("Appropriate for FastAPI SSR apps", 0),
            ("✓  app2 is a confidential client", 0),
        ],
        "Public Client",
        [
            ("SPA (React, Vue) or mobile app", 0),
            ("Cannot hold a secret — would be exposed", 0),
            ("Uses PKCE instead of client secret", 0),
            ("(Proof Key for Code Exchange)", 0),
            ("Required for client-side JavaScript apps", 0),
            ("✗  Not appropriate for app2", 0),
        ]
    )

    slide_bullets(prs, "App Registration: Step by Step", [
        ("1.   Log in to open.epic.com → My Apps → Create", 0),
        ("2.   Set Application Audience: Patients", 0),
        ("3.   Check Is Confidential Client: Yes", 0),
        ("4.   Set Redirect URI: http://localhost:8000/auth/callback", 0),
        ("5.   Select API Capabilities individually — do NOT select 'All FHIR R4 scopes'", 0),
        ("     'All scopes' generates a ~5000-char scope string that breaks the cookie", 1),
        ("6.   Submit and wait for Epic to provision (usually a few minutes)", 0),
        ("7.   Copy Client ID → set EPIC_NONPROD_CLIENT_ID in .env", 0),
        ("8.   Copy Client Secret immediately — the portal later shows only a hash", 0),
        ("", 0),
        ("JWK Set URLs: Leave blank (only needed for private_key_jwt auth — not used here)", 0),
        ("Persistent access: Not required (app2 only makes calls during active sessions)", 0),
    ])

    slide_bullets(prs, "API Capabilities: Observation is Category-Specific", [
        ("Most resources: one capability toggle covers the entire resource type", 0),
        ("Observation is different — Epic registers access per clinical category:", 0),
        ("", 0),
        ("Observation.Read — Vital Signs (R4)   →  grants category: vital-signs", 0),
        ("Observation.Read — Laboratory (R4)    →  grants category: laboratory", 0),
        ("Observation.Read — Social History (R4) →  grants category: survey", 0),
        ("", 0),
        ("Each enabled capability appears as a separate scope in the token response:", 0),
        ("patient/Observation.r?category=...observation-category|laboratory", 1),
        ("patient/Observation.r?category=...observation-category|vital-signs", 1),
        ("", 0),
        ("Epic enforces grants independently — a token with only 'laboratory'", 0),
        ("access returns 403 on a vital-signs request, even though it's the same resource", 0),
    ])

    slide_bullets(prs, "Sandbox vs Production Scope Behavior", [
        ("EPIC_SCOPE in .env:  'openid fhirUser'  (just two identity scopes)", 0),
        ("", 0),
        ("In Epic's SANDBOX:", 0),
        ("EPIC_SCOPE is a trigger to initiate the OAuth flow", 1),
        ("Epic grants ALL registered API capabilities regardless of what scopes were requested", 1),
        ("The app registration controls what you get — not the scope string", 1),
        ("Diagnosing 403s: check the portal registration, not EPIC_SCOPE", 1),
        ("", 0),
        ("In Epic's PRODUCTION environment:", 0),
        ("Behavior follows the SMART on FHIR specification strictly", 1),
        ("Must explicitly request each scope the app needs in EPIC_SCOPE", 1),
        ("Only scopes that were requested AND registered are granted", 1),
        ("'openid fhirUser' in production = no clinical data access", 1),
        ("", 0),
        ("Before deploying to production: enumerate all required scopes in EPIC_SCOPE", 0),
    ])


def gen_s6_oauth(prs):
    slide_section(prs, 6, "OAuth 2.0 Authorization Code Flow",
        "The five-step SMART on FHIR login sequence")

    slide_bullets(prs, "The Authorization Code Flow: Overview", [
        ("The Authorization Code flow is the standard for server-side web apps", 0),
        ("It keeps the access token off the browser's address bar and history", 0),
        ("", 0),
        ("Front-channel steps (browser redirects — address bar changes):", 0),
        ("Step 1  —  app2 redirects the browser to Epic's login page", 1),
        ("Step 2  —  user logs in and approves consent on Epic's page", 1),
        ("Step 3  —  Epic redirects browser back to app2 with a short-lived code", 1),
        ("", 0),
        ("Back-channel step (server-to-server — browser not involved):", 0),
        ("Step 4  —  app2 POSTs the code directly to Epic's token endpoint", 1),
        ("           Epic returns the access token — it never touches the browser", 1),
        ("", 0),
        ("Step 5  —  app2 stores the token server-side and sets a session cookie", 0),
    ])

    slide_bullets(prs, "Step 1: Building the Authorization URL", [
        ("GET /auth/login  — the entry point for the OAuth flow", 0),
        ("", 0),
        ("Generate a cryptographically random state value:", 0),
        ("state = secrets.token_urlsafe(32)  — 32 bytes → ~43-char URL-safe string", 1),
        ("Store in session: request.session['oauth_state'] = state", 1),
        ("", 0),
        ("Build the Epic authorization URL with these query parameters:", 0),
        ("response_type=code", 1),
        ("client_id=<EPIC_NONPROD_CLIENT_ID>", 1),
        ("scope=<EPIC_SCOPE (URL-encoded — forward slashes must be percent-encoded)>", 1),
        ("state=<random-state-value>", 1),
        ("redirect_uri=<APP_REDIRECT_URI>", 1),
        ("", 0),
        ("Respond with 307 Temporary Redirect to the Epic authorization URL", 0),
    ])

    slide_bullets(prs, "Step 2: State Parameter and CSRF Protection", [
        ("The state parameter is a one-time token binding the callback to the session", 0),
        ("Without it, the callback endpoint is vulnerable to login CSRF:", 0),
        ("", 0),
        ("The login CSRF attack:", 0),
        ("1. Attacker captures a valid authorization code from their own flow", 1),
        ("2. Tricks victim into visiting the callback URL with attacker's code", 1),
        ("3. Victim's browser completes the callback — token stored in victim's session", 1),
        ("4. Victim is now logged in as the attacker — or worse, data is cross-contaminated", 1),
        ("", 0),
        ("State prevents this:", 0),
        ("The callback checks state from Epic against state stored in the session", 1),
        ("Victim has no oauth_state in their session — check fails, 400 returned", 1),
        ("pop() is used (not get()) — makes state single-use, prevents replay", 1),
        ("", 0),
        ("In a FHIR app, a successful login CSRF could expose patient records", 0),
    ])

    slide_bullets(prs, "Step 3: Back-Channel Token Exchange", [
        ("GET /auth/callback  receives the authorization code from Epic", 0),
        ("Immediately POST to Epic's token endpoint  (back-channel — browser not involved):", 0),
        ("", 0),
        ("POST body (client_secret_post authentication):", 0),
        ("grant_type=authorization_code", 1),
        ("code=<authorization_code>  — expires in ~60 seconds, exchange immediately", 1),
        ("redirect_uri=<must match registration exactly>", 1),
        ("client_id=<EPIC_NONPROD_CLIENT_ID>", 1),
        ("client_secret=<EPIC_CLIENT_SECRET>", 1),
        ("", 0),
        ("Epic responds with:", 0),
        ("access_token, token_type, expires_in, scope", 1),
        ("id_token  (when openid scope is granted)", 1),
        ("patient  (the FHIR ID of the in-context patient — patient-scoped launches)", 1),
    ])

    slide_bullets(prs, "Step 4: Session Storage After Token Exchange", [
        ("Server-side token store (app.state.token_store — keyed by random session_id):", 0),
        ("access_token  — the JWT Bearer credential for FHIR API calls", 1),
        ("refresh_token — for obtaining new access tokens (not yet used in app2)", 1),
        ("scope         — granted scope string", 1),
        ("id_token      — OIDC identity JWT for the ID Token dialog", 1),
        ("patient       — Epic patient FHIR ID (required for MedicationRequest, Observation)", 1),
        ("", 0),
        ("Session cookie (signed, not encrypted — kept small):", 0),
        ("session_id         — random key into the server-side token store", 1),
        ("token_expires_at   — ISO 8601 expiry timestamp (checked before FHIR calls)", 1),
        ("scope              — stored for display on the home page", 1),
        ("", 0),
        ("Tokens are NEVER in the cookie — only a pointer to where they live server-side", 0),
    ])

    slide_bullets(prs, "The Meta-Refresh Trick (Why Not a 302 Redirect?)", [
        ("After storing the session, we need to send the browser to /", 0),
        ("A 302 redirect would be the natural choice — but it breaks in Chrome and Safari", 0),
        ("", 0),
        ("The problem:", 0),
        ("The callback is part of a cross-site redirect chain starting from Epic's domain", 1),
        ("Browsers treat SameSite=Lax cookies set on 302 responses mid-chain as ineligible", 1),
        ("The session cookie silently disappears — no error, just no session", 1),
        ("", 0),
        ("The fix — return a 200 with a client-side meta-refresh:", 0),
        ("<meta http-equiv='refresh' content='1;url=/'> ", 1),
        ("", 0),
        ("Why this works:", 0),
        ("The 200 response 'lands' the browser on our origin and stores the cookie", 1),
        ("The subsequent navigation to / is a same-site top-level GET — cookie is sent", 1),
    ])


def gen_s7_jwt(prs):
    slide_section(prs, 7, "JWTs and Token Management",
        "Access tokens, identity tokens, session cookies, and the token store")

    slide_bullets(prs, "What is a JWT?", [
        ("JSON Web Token — pronounced 'jot'", 0),
        ("A compact, URL-safe string carrying signed claims", 0),
        ("Used as both access tokens and ID tokens in SMART on FHIR", 0),
        ("", 0),
        ("Structure: three Base64URL-encoded segments joined by dots:", 0),
        ("<header>.<payload>.<signature>", 1),
        ("All three segments start with 'eyJ' when Base64URL-decoded from JSON", 1),
        ("", 0),
        ("Key property: the payload can be decoded without any key (it is not encrypted)", 0),
        ("The signature proves the token was issued by Epic and has not been altered", 1),
        ("Signature uses RS256: Epic's RSA private key + SHA-256", 1),
        ("", 0),
        ("Verification is the FHIR server's job — app2 forwards tokens, doesn't verify them", 0),
    ])

    slide_bullets(prs, "JWT Payload: Key Claims", [
        ("iss   — Issuer: Epic's OAuth base URL", 0),
        ("sub   — Subject: the authenticated user's Epic internal ID", 0),
        ("aud   — Audience: who the token is intended for (Epic FHIR base URL)", 0),
        ("client_id  — the registered client ID of this application", 0),
        ("iat   — Issued-at time (Unix timestamp)", 0),
        ("exp   — Expiration time (Unix timestamp) — Epic typically issues 1-hour tokens", 0),
        ("jti   — JWT ID: unique identifier for this specific token instance", 0),
        ("scope — Granted scopes (space-separated, Epic uses .r / .s abbreviations)", 0),
        ("fhirUser  — Full FHIR URL of the authenticated user's Patient/Practitioner resource", 0),
        ("patient   — Short-form patient context ID from the launch context", 0),
    ])

    slide_two_col(prs,
        "Access Token vs ID Token",
        "access_token  (Bearer credential)",
        [
            ("Authorization credential for FHIR API calls", 0),
            ("Intended for the resource server (Epic FHIR)", 0),
            ("Treat as opaque — forward unchanged", 0),
            ("Don't decode or build logic on its claims", 0),
            ("Only the FHIR server should interpret it", 0),
            ("Sent in Authorization: Bearer header", 0),
        ],
        "id_token  (OIDC identity assertion)",
        [
            ("Assertion about who authenticated", 0),
            ("Intended for the client app (app2)", 0),
            ("Decode to learn the user's identity", 0),
            ("sub, fhirUser, name, email...", 0),
            ("app2 decodes it for the ID Token dialog", 0),
            ("Never sent to the FHIR server", 0),
        ]
    )

    slide_bullets(prs, "The Server-Side Token Store", [
        ("app.state.token_store = {}  — in-memory Python dict, keyed by session_id", 0),
        ("Created during lifespan startup; shared across all requests", 0),
        ("", 0),
        ("Why not put the token in the session cookie?", 0),
        ("An Epic JWT access token is ~880 characters", 1),
        ("With refresh token + scope string, JSON-encoded and Base64-signed → easily > 4 KB", 1),
        ("Browsers silently discard cookies > 4096 bytes — session disappears with no error", 1),
        ("", 0),
        ("Solution: store only a short random session_id in the cookie:", 0),
        ("session_id = secrets.token_hex(16)  — 32 hex characters", 1),
        ("token_store[session_id] = { access_token, id_token, patient, ... }", 1),
        ("", 0),
        ("Limitation: in-memory — cleared on server restart → user must reconnect", 0),
        ("Production fix: replace with Redis or a database-backed store", 0),
    ])

    slide_bullets(prs, "The ID Token Dialog", [
        ("When openid scope is granted, Epic returns both access_token and id_token", 0),
        ("id_token is an OIDC JWT — contains identity claims about who authenticated", 0),
        ("", 0),
        ("Server-side decoding in the home route handler:", 0),
        ("_decode_jwt_payload() splits on '.', Base64URL-decodes the middle segment", 1),
        ("Returns a dict of claims — no signature verification (display only)", 1),
        ("Passes id_token_claims_json (pre-formatted) to the template", 1),
        ("", 0),
        ("Native HTML <dialog> element — no JavaScript library needed:", 0),
        ("showModal()  — opens as modal, dims backdrop, traps focus", 1),
        ("ESC key      — closes natively, no code required", 1),
        ("Backdrop click  — onclick checks event.target === this, calls close()", 1),
        ("", 0),
        ("| safe filter used for id_token_claims_json — trusted server-generated content", 0),
    ])


def gen_s8_fhir(prs):
    slide_section(prs, 8, "FHIR API Calls",
        "The six-step handler pattern and each endpoint in detail")

    slide_bullets(prs, "The Six-Step Handler Pattern", [
        ("Every FHIR route handler in pages.py follows the same six steps:", 0),
        ("", 0),
        ("Step 1  —  Read session_id from request.session", 0),
        ("          Return error fragment immediately if no active session", 1),
        ("Step 2  —  Look up access_token (and patient where required) from token_store", 0),
        ("          Return error fragment if missing (server restarted, session orphaned)", 1),
        ("Step 3  —  Check token expiry via token_expires_at in session cookie", 0),
        ("          Avoids sending a request Epic will reject with 401", 1),
        ("Step 4  —  Call Epic FHIR API with Authorization: Bearer + Accept: application/fhir+json", 0),
        ("          Wrapped in try/except httpx.TransportError for network failures", 1),
        ("Step 5  —  Handle non-2xx responses via _error_html (with response body as detail)", 0),
        ("Step 6  —  json.dumps(response.json(), indent=2) → HTMLResponse(<pre>...)</pre>)", 0),
    ])

    slide_bullets(prs, "GET /Patient: The 'Me' Context", [
        ("Patient is special — Epic automatically scopes it to the session patient", 0),
        ("No explicit patient ID search parameter is needed", 0),
        ("", 0),
        ("The FHIR call:", 0),
        ("GET /Patient  (no query params)", 1),
        ("Authorization: Bearer <access_token>", 1),
        ("Accept: application/fhir+json", 1),
        ("", 0),
        ("Epic's FHIR server infers the patient from the access token's claims", 0),
        ("The patient claim in the JWT identifies which patient's data to return", 1),
        ("", 0),
        ("Step 2 for /Patient only needs access_token — no patient ID lookup required", 0),
        ("Only the Patient resource type supports this me-context shortcut", 0),
    ])

    slide_bullets(prs, "GET /MedicationRequest: Patient-Scoped Search", [
        ("MedicationRequest has no 'me' context — requires an explicit patient parameter", 0),
        ("Without it, Epic rejects the request", 0),
        ("", 0),
        ("Where the patient ID comes from:", 0),
        ("Epic includes a 'patient' field in the token response for patient-scoped launches", 1),
        ("Stored in token_store['patient'] during /auth/callback", 1),
        ("Retrieved in step 2 of the handler: token_entry.get('patient', '')", 1),
        ("", 0),
        ("The FHIR call:", 0),
        ("GET /MedicationRequest?patient=<fhir_id>", 1),
        ("", 0),
        ("Response: FHIR Bundle where each entry is a MedicationRequest resource", 0),
        ("Each represents one prescription or medication order for the patient", 1),
        ("", 0),
        ("Required: patient/MedicationRequest.read scope + capability in app registration", 0),
    ])

    slide_bullets(prs, "The Observation Resource", [
        ("FHIR's catch-all for clinical measurements and findings", 0),
        ("One resource type covers an enormous range of clinical data:", 0),
        ("Lab results, vital signs, survey responses, smoking status, body measurements...", 1),
        ("", 0),
        ("Consistent structure regardless of measurement type:", 0),
        ("code      — what was measured (usually a LOINC code, e.g. 8867-4 = Heart rate)", 1),
        ("value[x]  — the result: a Quantity, a CodeableConcept, a string, or a boolean", 1),
        ("category  — the clinical domain (laboratory, vital-signs, survey...)", 1),
        ("subject   — reference to the patient", 1),
        ("effectiveDateTime  — when the observation was made", 1),
        ("", 0),
        ("Filtering by category is essential for retrieving a coherent subset", 0),
        ("Without category, you get everything mixed together in one large Bundle", 1),
    ])

    slide_bullets(prs, "Observation: The Category System", [
        ("Category system: http://terminology.hl7.org/CodeSystem/observation-category", 0),
        ("Pass only the code in the URL query string — not the full system URL", 0),
        ("", 0),
        ("Categories used in app2:", 0),
        ("laboratory  — CBC, metabolic panels, urinalysis, cultures, lipid panels", 1),
        ("vital-signs — blood pressure, heart rate, temperature, SpO2, height, weight, BMI", 1),
        ("", 0),
        ("Other standard categories:", 0),
        ("survey        — SDOH questionnaires, PHQ-9, GAD-7", 1),
        ("imaging       — radiology results", 1),
        ("social-history — smoking status, alcohol use", 1),
        ("", 0),
        ("category=vital-signs uses a hyphen — lowercase, exact match required", 0),
        ("The full param: ?patient=<id>&category=vital-signs", 0),
    ])

    slide_bullets(prs, "GET /Observation: Lab Reports + Vital Signs", [
        ("Both endpoints use the same /Observation resource — only the category differs:", 0),
        ("", 0),
        ("Lab reports:", 0),
        ("GET /Observation?patient=<id>&category=laboratory", 1),
        ("Requires: Observation.Read — Laboratory (R4) capability in app registration", 1),
        ("Scope granted: patient/Observation.r?category=...observation-category|laboratory", 1),
        ("", 0),
        ("Vital signs:", 0),
        ("GET /Observation?patient=<id>&category=vital-signs", 1),
        ("Requires: Observation.Read — Vital Signs (R4) capability in app registration", 1),
        ("Scope granted: patient/Observation.r?category=...observation-category|vital-signs", 1),
        ("", 0),
        ("Both follow the identical six-step handler pattern", 0),
        ("The patient ID comes from token_store['patient'] — same as MedicationRequest", 0),
        ("Empty bundle (total: 0) is a valid 200 response — not an error", 0),
    ])

    slide_bullets(prs, "Observation: Result Volume and Pagination", [
        ("Observation bundles can be significantly larger than MedicationRequest bundles", 0),
        ("A patient with years of history may have hundreds of lab results", 0),
        ("", 0),
        ("To limit results, add _count to the query params:", 0),
        ("params={'patient': patient_id, 'category': 'laboratory', '_count': 20}", 1),
        ("", 0),
        ("When more results exist, Epic includes a Bundle.link with relation='next':", 0),
        ("{ 'relation': 'next', 'url': 'https://fhir.epic.com/...?page=2' }", 1),
        ("Follow the next URL to retrieve the next page of results", 1),
        ("", 0),
        ("app2 does not currently implement _count or pagination", 0),
        ("Sufficient for sandbox test patients with limited histories", 1),
        ("Add _count before deploying against real patients with extended histories", 1),
    ])


def gen_s9_patterns(prs):
    slide_section(prs, 9, "Key Design Patterns",
        "Lifespan, error handling, token security, and production readiness")

    slide_bullets(prs, "Lifespan and app.state", [
        ("FastAPI's lifespan context manager runs once at startup and once at shutdown", 0),
        ("The correct place to create and tear down application-level resources", 0),
        ("", 0),
        ("app2's lifespan creates three resources:", 0),
        ("app.state.http_client = httpx.AsyncClient()  — shared connection pool", 1),
        ("app.state.templates   = Jinja2Templates(...)  — pre-loaded template environment", 1),
        ("app.state.token_store = {}                   — in-memory session token store", 1),
        ("", 0),
        ("Shutdown: await app.state.http_client.aclose()  — clean TCP connection drain", 0),
        ("", 0),
        ("Why not module-level globals?", 0),
        ("Event loop doesn't exist at import time — AsyncClient must be created after startup", 1),
        ("Module globals create invisible dependencies; app.state makes them explicit", 1),
        ("Lifespan resources are scoped to the app instance — testable in isolation", 1),
    ])

    slide_bullets(prs, "Token Security: Defense in Depth", [
        ("Access tokens are Bearer credentials — whoever holds one can make FHIR calls", 0),
        ("Every design decision in app2 minimizes exposure:", 0),
        ("", 0),
        ("Back-channel token exchange:", 0),
        ("The access token is retrieved server-to-server — never via the browser", 1),
        ("", 0),
        ("Server-side token store:", 0),
        ("Tokens stored in app.state.token_store — never in the session cookie", 1),
        ("Even if an attacker reads the cookie, they only get a session_id and metadata", 1),
        ("", 0),
        ("Minimal session cookie payload:", 0),
        ("Cookie is signed (itsdangerous) but not encrypted — payload is readable", 1),
        ("Only session_id, token_expires_at, and scope go in the cookie", 1),
        ("No tokens, no patient IDs, no PII", 1),
        ("", 0),
        ("Bearer tokens are forwarded opaquely — never decoded, never logged", 0),
    ])

    slide_bullets(prs, "Production Readiness Checklist", [
        ("SessionMiddleware: set https_only=True (requires HTTPS/TLS termination)", 0),
        ("Reverse proxy: run uvicorn behind nginx (handles TLS, static files, keepalive)", 0),
        ("Process manager: systemd or supervisor — not --reload (dev only)", 0),
        ("", 0),
        ("Token store: replace app.state.token_store with Redis", 0),
        ("In-memory store is cleared on restart — users must re-authenticate", 1),
        ("Redis survives restarts and supports multiple worker processes", 1),
        ("", 0),
        ("EPIC_SCOPE: enumerate all required resource scopes explicitly for production", 0),
        ("Sandbox grants all registered capabilities; production grants only what's requested", 1),
        ("", 0),
        ("httpx timeouts: add Timeout(connect=5.0, read=30.0) to AsyncClient", 0),
        ("Without timeouts, a hung Epic endpoint holds the handler open indefinitely", 1),
        ("", 0),
        ("APP_REDIRECT_URI: update to the production HTTPS callback URL", 0),
    ])


def gen_s10_build(prs):
    slide_section(prs, 10, "Building Your Own",
        "Getting started, diagnosing errors, and extending the app")

    slide_bullets(prs, "Getting Started: The Critical Path", [
        ("1.  Set up Python 3.12 environment and install dependencies", 0),
        ("    python3.12 -m venv .venv && source .venv/bin/activate", 1),
        ("    pip install -r requirements.txt", 1),
        ("2.  Register a confidential client app at open.epic.com", 0),
        ("    Enable only the API capabilities you need (not 'All FHIR R4 scopes')", 1),
        ("3.  Create .env with CLIENT_ID, CLIENT_SECRET, SESSION_SECRET_KEY", 0),
        ("    Generate SESSION_SECRET_KEY: python -c \"import secrets; print(secrets.token_hex(32))\"", 1),
        ("4.  Create __init__.py in the app directory (required for dot-notation uvicorn)", 0),
        ("5.  Build main.py: FastAPI + lifespan + SessionMiddleware + include_router", 0),
        ("6.  Build routers/auth.py: /auth/login, /auth/callback, /auth/logout", 0),
        ("7.  Build routers/pages.py: home, portal, FHIR endpoints", 0),
        ("8.  Run: uvicorn app2.main:app --reload --port 8000", 0),
    ])

    slide_bullets(prs, "Diagnosing Common Errors", [
        ("403 on FHIR call:", 0),
        ("Check the Scope row on the home page — is the required scope present?", 1),
        ("If not: enable the capability in the Epic portal → log out → reconnect", 1),
        ("", 0),
        ("Session disappears after OAuth callback:", 0),
        ("SameSite/cross-site redirect issue — use meta-refresh (200) not 302 redirect", 1),
        ("", 0),
        ("400 State mismatch:", 0),
        ("Session not persisting between /auth/login and /auth/callback", 1),
        ("Check SESSION_SECRET_KEY is set and consistent", 1),
        ("", 0),
        ("Token store empty after clicking FHIR button:", 0),
        ("Server restarted — in-memory store cleared. Log out and reconnect.", 1),
        ("", 0),
        ("ReadError on FHIR call:", 0),
        ("Stale pooled TCP connection — retry once, usually succeeds immediately", 1),
    ])

    slide_bullets(prs, "Extending the App: What's Next", [
        ("Add _count pagination for Observation endpoints", 0),
        ("Implement Bundle.link traversal for large result sets", 1),
        ("", 0),
        ("Refactor duplicate Observation handlers into a shared helper", 0),
        ("_fhir_get_observation(request, category) called by two thin route handlers", 1),
        ("", 0),
        ("Add more FHIR resource endpoints:", 0),
        ("GET /Condition?patient=<id>&category=problem-list-item  — diagnoses", 1),
        ("GET /AllergyIntolerance?patient=<id>  — allergies", 1),
        ("GET /Encounter?patient=<id>  — past visits", 1),
        ("GET /Immunization?patient=<id>  — vaccination history", 1),
        ("", 0),
        ("Production upgrades:", 0),
        ("Replace in-memory token store with Redis", 1),
        ("Add HTTPS + nginx reverse proxy", 1),
        ("Deploy to OCI (Oracle Cloud Infrastructure) VM — lift and shift", 1),
    ])


def gen_s11_close(prs):
    slide_section(prs, 11, "Key Resources", "")

    slide_bullets(prs, "Reference Links", [
        ("SMART on FHIR specification:", 0),
        ("smart.hl7.org", 1),
        ("", 0),
        ("FHIR R4 specification:", 0),
        ("hl7.org/fhir/R4", 1),
        ("", 0),
        ("Epic developer portal and sandbox:", 0),
        ("open.epic.com  (app registration, API capability documentation)", 1),
        ("", 0),
        ("FastAPI documentation:", 0),
        ("fastapi.tiangolo.com", 1),
        ("", 0),
        ("HTMX documentation:", 0),
        ("htmx.org", 1),
        ("", 0),
        ("HTTPX documentation:", 0),
        ("python-httpx.org", 1),
        ("", 0),
        ("This repository's companion developer guide:", 0),
        ("docs/app2-guide.md  — every implementation step documented in full detail", 1),
    ])

    slide_closing(prs,
        "Thank You",
        "Questions?  ·  Refer to docs/app2-guide.md for full implementation details")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    gen_title(prs)
    gen_s1_intro(prs)
    gen_s2_fhir(prs)
    gen_s3_arch(prs)
    gen_s4_stack(prs)
    gen_s5_reg(prs)
    gen_s6_oauth(prs)
    gen_s7_jwt(prs)
    gen_s8_fhir(prs)
    gen_s9_patterns(prs)
    gen_s10_build(prs)
    gen_s11_close(prs)

    out = "docs/app2-tutorial.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")

if __name__ == "__main__":
    main()
